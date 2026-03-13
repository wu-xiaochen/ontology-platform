#!/usr/bin/env python3
"""生成 OpenClaw 2026 应用分析演示文稿"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色主题 - Ocean Gradient
PRIMARY = RGBColor(6, 90, 130)      # deep blue
SECONDARY = RGBColor(28, 114, 147)   # teal
ACCENT = RGBColor(33, 41, 92)        # midnight
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(242, 242, 242)

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景形状
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = SECONDARY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, accent_color=PRIMARY):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = accent_color
    
    # 要点列表
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
    
    return slide

def add_stat_slide(prs, title, stats):
    """添加数据展示页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 统计数据
    cols = len(stats)
    col_width = Inches(12) / cols
    
    for i, stat in enumerate(stats):
        x = Inches(0.5) + i * col_width
        # 数字
        num_box = slide.shapes.add_textbox(x, Inches(2), col_width, Inches(2))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat["value"]
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = SECONDARY
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(x, Inches(4.2), col_width, Inches(1))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat["label"]
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# ===== 创建幻灯片 =====

# 1. 标题页
add_title_slide(prs, 
    "OpenClaw 2026 应用分析",
    "AI Agent 平台的演进与未来趋势")

# 2. 目录
add_content_slide(prs, "目录", [
    "OpenClaw 平台概述",
    "核心功能与技术架构",
    "2026 年主要应用场景",
    "社区生态与插件系统",
    "未来展望与机会"
])

# 3. 平台概述
add_content_slide(prs, "OpenClaw 平台概述", [
    "开源 AI Agent 编排框架",
    "支持多种大语言模型集成 (Claude, GPT, Gemini, MiniMax等)",
    "灵活的插件系统与技能市场",
    "跨平台部署能力 (本地、云端、边缘设备)",
    "强大的工作流自动化能力"
])

# 4. 核心架构
add_content_slide(prs, "核心功能与技术架构", [
    "多模型路由与负载均衡",
    "可扩展的技能 (Skills) 系统",
    "实时消息通道集成 (Telegram, Discord, Slack等)",
    "Memory 与知识图谱持久化",
    "安全审计与权限管理"
])

# 5. 统计数据
add_stat_slide(prs, "2026 年关键数据", [
    {"value": "18,987+", "label": "Skills 数量"},
    {"value": "88.7K+", "label": "GitHub Stars"},
    {"value": "50+", "label": "集成平台"},
    {"value": "100K+", "label": "活跃用户"}
])

# 6. 应用场景 1
add_content_slide(prs, "2026 年主要应用场景", [
    "🏠 智能家居与 IoT 设备控制",
    "💼 企业工作流自动化",
    "📱 跨平台消息与通讯管理",
    "🧪 开发者工具与代码辅助",
    "📊 数据分析与报告生成"
])

# 7. 应用场景 2
add_content_slide(prs, "个人 AI 助手应用", [
    "日程管理与智能提醒",
    "邮件自动分类与回复建议",
    "文件整理与知识管理",
    "多语言翻译与内容创作",
    "个性化学习与技能提升"
])

# 8. 社区生态
add_content_slide(prs, "社区生态与插件系统", [
    "ClawHub 技能市场 - 一键安装社区贡献的技能",
    "开放插件接口 - 轻松集成第三方服务",
    "开发者工具链 - VS Code, CLI, SDK",
    "丰富的预置模板 - 开箱即用的解决方案",
    "活跃的 Discord 社区 - 经验交流与支持"
])

# 9. 未来展望
add_content_slide(prs, "未来展望与机会", [
    "多模态能力增强 - 图像、语音、视频处理",
    "更强的自主决策能力 - 复杂任务规划",
    "边缘计算优化 - 低延迟本地部署",
    "企业级安全增强 - 合规与审计",
    "AI Agent 市场 - 专业化垂直解决方案"
])

# 10. 结束页
add_title_slide(prs, 
    "谢谢！",
    "让 AI 成为你的第二大脑")

# 保存文件
output_path = os.path.expanduser("~/Desktop/OpenClaw_2026_Analysis.pptx")
prs.save(output_path)
print(f"演示文稿已保存到: {output_path}")
